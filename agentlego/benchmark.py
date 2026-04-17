import csv
import importlib
import json
import os
import re
from pathlib import Path
from typing import List, Optional, Sequence, Tuple, Union

from agentlego.types import ImageIO, VideoIO, AudioIO, Annotated, File, Info
from agentlego.utils import load_or_build_object, require, parse_multi_float
from agentlego.tools.base import BaseTool
from agentlego.tools.calculator.python_calculator import Calculator
from agentlego.tools.search.google import GoogleSearch
from agentlego.tools.ocr.ocr import OCR
from agentlego.tools.ocr.math_ocr import MathOCR
from agentlego.tools.image_editing.draw_box import DrawBox
from agentlego.tools.image_editing.add_text import AddText
from agentlego.tools.python_interpreter.plot import Plot
from agentlego.tools.python_interpreter.solver import Solver
from agentlego.tools.speech_text.speech_to_text import SpeechToText

class _LazyModule:
    """A tiny lazy import proxy.

    This keeps module import side-effects out of the global import path so that
    older/minimal tool deployments can import this file without installing every
    optional dependency.
    """

    def __init__(self, module_name: str, pip_name: Optional[str] = None):
        self._module_name = module_name
        self._pip_name = pip_name or module_name
        self._module = None

    def _load(self):
        if self._module is None:
            try:
                self._module = importlib.import_module(self._module_name)
            except ModuleNotFoundError as e:
                raise ModuleNotFoundError(
                    f"Optional dependency '{self._pip_name}' is required for this tool. "
                    f"Please install it to use this functionality."
                ) from e
        return self._module

    def __getattr__(self, name: str):
        return getattr(self._load(), name)


# Common optional third-party dependencies used by some tools below.
# They are imported lazily so importing this module does not require them.
cv2 = _LazyModule('cv2', 'opencv-python')
np = _LazyModule('numpy', 'numpy')
openpyxl = _LazyModule('openpyxl', 'openpyxl')
nr = _LazyModule('noisereduce', 'noisereduce')
sf = _LazyModule('soundfile', 'soundfile')
pd = _LazyModule('pandas', 'pandas')
PyPDF2 = _LazyModule('PyPDF2', 'PyPDF2')
PILImage = _LazyModule('PIL.Image', 'Pillow')

class QwenVLInferencer:

    def __init__(self,
                 model='Qwen/Qwen-VL-Chat',
                 revision='f57cfbd358cb56b710d963669ad1bcfb44cdcdd8',
                 fp16=False,
                 device=None):
        from transformers import AutoModelForCausalLM, AutoTokenizer
        self.tokenizer = AutoTokenizer.from_pretrained(
            model, trust_remote_code=True, revision=revision)

        self.model = AutoModelForCausalLM.from_pretrained(
            model,
            device_map=device or 'auto',
            trust_remote_code=True,
            revision=revision,
            fp16=fp16).eval()

    def __call__(self, image: ImageIO, text: str):
        query = self.tokenizer.from_list_format(
            [dict(image=image.to_path()), dict(text=text)])
        response, _ = self.model.chat(self.tokenizer, query=query, history=None)
        return response

    def fetch_all_box_with_ref(self, text) -> List[dict]:
        items = self.tokenizer._fetch_all_box_with_ref(text)
        return items



class ImageDescription(BaseTool):
    default_desc = ('A useful tool that returns a brief '
                    'description of the input image.')

    @require('mmpretrain')
    def __init__(self,
                 model: str = 'llava-7b-v1.5_vqa',
                 device: str = 'cuda',
                 toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.model = model
        self.device = device

    def setup(self):
        from mmengine.registry import DefaultScope
        from mmpretrain.apis import VisualQuestionAnsweringInferencer
        with DefaultScope.overwrite_default_scope('mmpretrain'):
            self._inferencer = load_or_build_object(
                VisualQuestionAnsweringInferencer,
                model=self.model,
                device=self.device,
            )

    def apply(self, image: ImageIO) -> str:
        image = image.to_array()[:, :, ::-1]
        return self._inferencer(image, 'Describe the image in detail')[0]['pred_answer']


class CountGivenObject(BaseTool):
    default_desc = 'The tool can count the number of a certain object in the image.'

    @require('mmpretrain')
    def __init__(self,
                 model: str = 'llava-7b-v1.5_vqa',
                 device: str = 'cuda',
                 toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.model = model
        self.device = device

    def setup(self):
        from mmengine.registry import DefaultScope
        from mmpretrain.apis import VisualQuestionAnsweringInferencer
        with DefaultScope.overwrite_default_scope('mmpretrain'):
            self._inferencer = load_or_build_object(
                VisualQuestionAnsweringInferencer,
                model=self.model,
                device=self.device,
            )

    def apply(
        self,
        image: ImageIO,
        text: Annotated[str, Info('The object description in English.')],
        bbox: Annotated[Optional[str],
                        Info('The bbox coordinate in the format of `(x1, y1, x2, y2)`')] = None,
    ) -> int:
        import re
        if bbox is None:
            image = image.to_array()[:, :, ::-1]
        else:
            x1, y1, x2, y2 = (int(item) for item in parse_multi_float(bbox))
            image = image.to_array()[y1:y2, x1:x2, ::-1]
        res = self._inferencer(image, f'How many {text} are in the image? Reply a digit')[0]['pred_answer']
        res = re.findall(r'\d+', res)
        if len(res) > 0:
            return int(res[0])
        else:
            return 0


class RegionAttributeDescription(BaseTool):
    default_desc = 'Describe the attribute of a region of the input image.'

    @require('mmpretrain')
    def __init__(self,
                 model: str = 'llava-7b-v1.5_vqa',
                 device: str = 'cuda',
                 toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.model = model
        self.device = device


    def setup(self):
        from mmengine.registry import DefaultScope
        from mmpretrain.apis import VisualQuestionAnsweringInferencer
        with DefaultScope.overwrite_default_scope('mmpretrain'):
            self._inferencer = load_or_build_object(
                VisualQuestionAnsweringInferencer,
                model=self.model,
                device=self.device,
            )

    def apply(
        self,
        image: ImageIO,
        bbox: Annotated[str,
                        Info('The bbox coordinate in the format of `(x1, y1, x2, y2)`')],
        attribute: Annotated[str, Info('The attribute to describe')],
    ) -> str:
        x1, y1, x2, y2 = (int(item) for item in parse_multi_float(bbox))
        cropped_image = image.to_array()[y1:y2, x1:x2, ::-1]
        return self._inferencer(cropped_image, f'Describe {attribute} on the image in detail')[0]['pred_answer']



class AudioClipTool(BaseTool):
    """A tool to extract a clip from an audio file between specified start and end times.

    Args:
        start_time (float): The start time of the audio clip in seconds.
        end_time (float): The end time of the audio clip in seconds.
        output_format (str): The format of the output audio (e.g., 'mp3', 'wav'). Defaults to 'mp3'.
        toolmeta (None | dict | ToolMeta): Additional metadata for the tool.
    """

    default_desc = "This tool extracts a clip from an audio file between specified start and end times."

    @require('pydub')
    def __init__(self,
                 start_time: float = 0.0,
                 end_time: float = 10.0,
                 output_format: str = 'mp3',
                 toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.start_time = start_time
        self.end_time = end_time
        self.output_format = output_format

    def setup(self):
        # No additional setup required.
        pass

    def apply(self,
              audio: str,
              start_time: Optional[float] = None,
              end_time: Optional[float] = None,
              filename: str = None,
              output_dir: str = None
              ) -> Annotated[AudioIO, Info("Extracted audio clip")]:
        """
        Extract a clip from `audio`.

        start_time and end_time (in seconds) may be provided here; if omitted,
        the instance defaults (from the constructor) are used.
        """
        # Input is expected to be a file path string. For robustness, also
        # accept AudioIO-like objects that expose `.to_path()`.
        audio_path = audio.to_path() if hasattr(audio, 'to_path') else str(audio)

        # Determine start/end seconds (use call-time args if provided)
        st = self.start_time if start_time is None else float(start_time)
        ed = self.end_time if end_time is None else float(end_time)

        # Basic validation
        if st is None or ed is None:
            raise ValueError("start_time and end_time must be provided either in constructor or apply() call")
        if st < 0 or ed <= st:
            raise ValueError("Invalid start_time/end_time: must be non-negative and end_time > start_time")

        # Convert to milliseconds
        start_ms = int(st * 1000)
        end_ms = int(ed * 1000)

        # Load audio and slice
        from pydub import AudioSegment
        audio_segment = AudioSegment.from_file(audio_path)
        clip = audio_segment[start_ms:end_ms]

        # Prepare output directory
        import os
        if output_dir is None:
            output_dir = DEFAULT_TOOLS_OUTPUT_DIR
        os.makedirs(output_dir, exist_ok=True)

        # Normalize extension and build filename
        ext = str(self.output_format).lstrip('.')
        if filename is None:
            # Create a safe filename (keep the real extension!)
            safe_st = f"{st:.3f}".replace('.', '_')
            safe_ed = f"{ed:.3f}".replace('.', '_')
            filename = f"audio_clip_{safe_st}_{safe_ed}.{ext}"
        else:
            if not filename.lower().endswith(f".{ext}"):
                filename = f"{filename}.{ext}"

        output_path = os.path.join(output_dir, filename)
        # Use normalized format name to avoid leading-dot issues.
        clip.export(output_path, format=ext)
        # Return an AudioIO object instead of a dict for consistency with other audio tools
        return AudioIO(output_path)

class TextEmotionRecognize(BaseTool):
    """
    A tool to recognize the dominant emotion in a given text.

    Args:
        model (str): The Hugging Face model to use for emotion recognition.
                     Defaults to "j-hartmann/emotion-english-distilroberta-base".
        device (str | bool | int): The device to load the model. For example, "cuda" for GPU or False for CPU.
                                   Defaults to True, which means automatically select device.
        toolmeta (None | dict | ToolMeta): Additional info for the tool. Defaults to None.
    """

    default_desc = 'This tool can recognize the dominant emotion in the input text.'

    @require('transformers')
    def __init__(self, model: str = "distilbert-base-uncased-finetuned-sst-2-english", device: Union[str, bool, int] = True, toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.model_name = model
        self.device = device

    def setup(self):
        from transformers import pipeline
        # Determine device id: use GPU (device=0) if requested, otherwise CPU (-1)
        device_id = -1
        if isinstance(self.device, str) and self.device.lower() == "cuda":
            device_id = 0
        elif isinstance(self.device, bool):
            device_id = 0 if self.device else -1
        elif isinstance(self.device, int):
            device_id = self.device

        # Initialize the Hugging Face pipeline for emotion recognition.
        self._emotion_pipeline = pipeline(
            "text-classification",
            model=self.model_name,
            return_all_scores=True,
            device=device_id
        )

    def apply(
        self,
        text: str,
    ) -> Annotated[str,
                   Info('The recognized emotion label and its score from the input text.')]:
        # Run the emotion recognition pipeline on the input text.
        results = self._emotion_pipeline(text)
        if not results or not results[0]:
            return "unknown (0.0)"
        
        # Select the emotion with the highest score.
        best_emotion = max(results[0], key=lambda x: x["score"])
        return f"{best_emotion['label']} ({best_emotion['score']:.4f})"


class CsvFileGenerator(BaseTool):
    default_desc = 'Generates a CSV file with the provided data.'

    def apply(self, filename: str, data_str: str, output_dir: str = None) -> Annotated[File, Info('Generated CSV file.', filetype='text/csv')]:
        import os
        try:
            data = json.loads(data_str)
            if output_dir is None:
                output_dir = DEFAULT_TOOLS_OUTPUT_DIR
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, filename + ".csv")
            with open(output_path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                writer.writerows(data)
            return File(output_path, filetype='text/csv')
        except json.JSONDecodeError:
            return 'Invalid data format: data should be a JSON-formatted string representing a list of lists.'
        except Exception as e:
            return f'An error occurred: {str(e)}'



class DocxFileGenerator(BaseTool):
    default_desc = (
        'A Word (.docx) document generator supporting Chinese characters and image embedding. '
        'Use placeholders like {{image_1}}, {{image_2}} in the "content" parameter to specify '
        'exactly where images should be inserted. The "image_paths" parameter must be a '
        'comma-separated string of file paths (e.g., "path/1.jpg, path/2.png") that matches '
        'the numerical order of the placeholders. The tool handles font encoding and image scaling.'
    )

    def apply(
        self, 
        filename: str, 
        content: str, 
        image_paths: str = "", # 接收逗号分隔字符串以适配框架
        output_dir: str = None,
        render_markdown: bool = False,
    ) -> Annotated[File, Info('Generated DOCX file.', filetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')]:
        try:
            from docx import Document
            from docx.oxml.ns import qn
            from docx.shared import Inches, Pt

            if output_dir is None:
                output_dir = DEFAULT_TOOLS_OUTPUT_DIR
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, filename + ".docx")

            doc = Document()
            
            # --- 1. 设置全局中文字体 ---
            # 这里的 'SimSun' 或 '宋体' 取决于你系统中安装的字体名称
            font_name = 'SimSun' 
            
            # 设置正文默认字体
            style = doc.styles['Normal']
            style.font.name = font_name
            style.font.size = Pt(12)
            # 必须设置这行，否则中文无法显示为指定字体
            style._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

            # --- 2. 处理图片路径列表 ---
            path_list = [p.strip() for p in image_paths.split(',') if p.strip()]

            # --- 3. 解析内容并插入文本/图片 ---
            if render_markdown:
                self._render_markdown_to_docx(doc, content, path_list, font_name)
            else:
                # 保持向后兼容的占位符方案
                parts = re.split(r'(\{\{image_\d+\}\})', content)

                for part in parts:
                    placeholder_match = re.match(r'\{\{image_(\d+)\}\}', part)
                    
                    if placeholder_match:
                        idx = int(placeholder_match.group(1)) - 1
                        if 0 <= idx < len(path_list):
                            img_path = path_list[idx]
                            if os.path.exists(img_path):
                                doc.add_picture(img_path, width=Inches(5))
                            else:
                                doc.add_paragraph(f"[图片文件未找到: {img_path}]")
                        else:
                            doc.add_paragraph(f"[缺失图片路径: {part}]")
                    else:
                        if part.strip():
                            for para in part.split('\n\n'):
                                p = doc.add_paragraph(para)
                                for run in p.runs:
                                    run.font.name = font_name
                                    run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

            doc.save(output_path)
            return File(output_path)

        except Exception as e:
            return f'An error occurred: {str(e)}'

    def _apply_inline_md_runs(self, paragraph, text, font_name):
        """将简单的 Markdown 内联样式渲染为 runs: **加粗**、*斜体*、`代码`、[链接](url)。"""
        from docx.oxml.ns import qn
        # 处理超链接 -> 以 "text (url)" 形式保留
        def repl_link(m):
            label, url = m.group(1), m.group(2)
            return f"{label} ({url})"

        text = re.sub(r"\[([^\]]+)\]\(([^\)]+)\)", repl_link, text)

        # 逐步拆分处理 `code`
        tokens = re.split(r"(`[^`]+`)", text)
        for tok in tokens:
            if not tok:
                continue
            if tok.startswith('`') and tok.endswith('`'):
                run = paragraph.add_run(tok[1:-1])
                run.font.name = 'Courier New'
                run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
                continue
            # 处理 **bold** 和 *italic*
            pos = 0
            pattern = re.compile(r"(\*\*[^*]+\*\*|\*[^*]+\*)")
            for m in pattern.finditer(tok):
                if m.start() > pos:
                    r = paragraph.add_run(tok[pos:m.start()])
                    r.font.name = font_name
                    r._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
                frag = m.group(0)
                if frag.startswith('**'):
                    r = paragraph.add_run(frag[2:-2])
                    r.bold = True
                    r.font.name = font_name
                    r._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
                else:
                    r = paragraph.add_run(frag[1:-1])
                    r.italic = True
                    r.font.name = font_name
                    r._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
                pos = m.end()
            if pos < len(tok):
                r = paragraph.add_run(tok[pos:])
                r.font.name = font_name
                r._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

    def _render_markdown_to_docx(self, doc, md_text: str, path_list: List[str], font_name: str):
        """轻量级 Markdown -> DOCX 渲染，支持：标题、段落、无序/有序列表、代码块、内联样式与图片。"""
        from docx.oxml.ns import qn
        from docx.shared import Inches
        lines = md_text.splitlines()
        i = 0
        while i < len(lines):
            line = lines[i]
            if not line.strip():
                i += 1
                continue

            # 代码块 ```
            if line.strip().startswith('```'):
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith('```'):
                    code_lines.append(lines[i])
                    i += 1
                # 跳过结束 ``` 行
                i += 1
                p = doc.add_paragraph('\n'.join(code_lines))
                for run in p.runs:
                    run.font.name = 'Courier New'
                    run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
                continue

            # 标题 #{1..6}
            m = re.match(r"^(#{1,6})\s+(.*)$", line)
            if m:
                level = min(len(m.group(1)), 6)
                text = m.group(2)
                h = doc.add_heading(level=level)
                self._apply_inline_md_runs(h, text, font_name)
                i += 1
                continue

            # 无序/有序列表
            if re.match(r"^\s*([*+-])\s+", line) or re.match(r"^\s*\d+\.\s+", line):
                is_ordered = bool(re.match(r"^\s*\d+\.\s+", line))
                while i < len(lines) and (re.match(r"^\s*([*+-])\s+", lines[i]) or re.match(r"^\s*\d+\.\s+", lines[i])):
                    item_text = re.sub(r"^\s*([*+-]|\d+\.)\s+", "", lines[i]).strip()
                    p = doc.add_paragraph(style='List Number' if is_ordered else 'List Bullet')
                    self._apply_inline_md_runs(p, item_text, font_name)
                    i += 1
                continue

            # 图片：![](path) 或 占位符 {{image_n}}
            img_md = re.match(r"^!\[[^\]]*\]\(([^\)]+)\)", line.strip())
            ph_md = re.match(r"^\{\{image_(\d+)\}\}$", line.strip())
            if img_md:
                img_path = img_md.group(1)
                if os.path.exists(img_path):
                    doc.add_picture(img_path, width=Inches(5))
                else:
                    doc.add_paragraph(f"[图片文件未找到: {img_path}]")
                i += 1
                continue
            if ph_md:
                idx = int(ph_md.group(1)) - 1
                if 0 <= idx < len(path_list):
                    img_path = path_list[idx]
                    if os.path.exists(img_path):
                        doc.add_picture(img_path, width=Inches(5))
                    else:
                        doc.add_paragraph(f"[图片文件未找到: {img_path}]")
                else:
                    doc.add_paragraph(f"[缺失图片路径: {{image_{idx+1}}}]")
                i += 1
                continue

            # 普通段落，支持行内样式和占位符穿插
            # 将行内 {{image_n}} 切分
            parts = re.split(r'(\{\{image_\d+\}\})', line)
            if len(parts) == 1:
                p = doc.add_paragraph()
                self._apply_inline_md_runs(p, line, font_name)
            else:
                p = doc.add_paragraph()
                for part in parts:
                    ph = re.match(r'\{\{image_(\d+)\}\}', part)
                    if ph:
                        idx = int(ph.group(1)) - 1
                        if 0 <= idx < len(path_list):
                            img_path = path_list[idx]
                            if os.path.exists(img_path):
                                doc.add_picture(img_path, width=Inches(5))
                            else:
                                doc.add_paragraph(f"[图片文件未找到: {img_path}]")
                        else:
                            doc.add_paragraph(f"[缺失图片路径: {part}]")
                    else:
                        self._apply_inline_md_runs(p, part, font_name)
            i += 1


# class DocxFileGenerator(BaseTool):
#     default_desc = 'Generates a .docx file with the provided content.'

#     def apply(self, filename: str, content: str, output_dir: str = None) -> Annotated[File, Info('Generated DOCX file.', filetype='office/docx')]:
#         import os
#         try:
#             if output_dir is None:
#                 output_dir = os.path.join(os.getcwd(), "tools_output")
#             os.makedirs(output_dir, exist_ok=True)
#             output_path = os.path.join(output_dir, filename + ".docx")
#             doc = Document()
#             doc.add_paragraph(content)
#             doc.save(output_path)
#             return File(output_path, filetype='office/docx')
#         except Exception as e:
#             return f'An error occurred: {str(e)}'



class PdfFileGenerator(BaseTool):
    default_desc = (
        'A PDF generator supporting Chinese and image embedding. Use {{image_n}} placeholders '
        'in the "content" to specify where images should appear. The "image_paths" parameter '
        'must be a comma-separated string of file paths matching the placeholder sequence.'
    )

    def apply(
        self, 
        filename: str, 
        content: str, 
        image_paths: str = "", # 接收逗号分隔的字符串路径
        output_dir: str = None,
        render_markdown: bool = False,
    ) -> Annotated[File, Info('Generated PDF file.', filetype='application/pdf')]:
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import inch
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            from reportlab.platypus import (
                SimpleDocTemplate,
                Paragraph,
                Image as RLImage,
                Spacer,
                ListFlowable,
                ListItem,
                Preformatted,
            )

            # 确定输出目录
            if output_dir is None:
                output_dir = DEFAULT_TOOLS_OUTPUT_DIR
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, filename + ".pdf")

            # --- 注册中文字体 ---
            # 默认字体名
            font_name = 'Helvetica'
            # 常见路径参考: Windows 为 "simsun.ttc", Linux 通常需要指向具体路径如 "/usr/share/fonts/..."
            font_path = "/home/zhenjue2/GTA-v2/agentlego/simsun.ttc" 
            
            try:
                # 只有在还没注册时才注册，避免重复注册报错
                if 'SimSun' not in pdfmetrics.getRegisteredFontNames():
                    pdfmetrics.registerFont(TTFont('SimSun', font_path))
                font_name = 'SimSun'
            except Exception as e:
                # 如果没找到指定字体，回退到系统默认
                print(f"Warning: Font file not found at {font_path}, fallback to Helvetica. Error: {str(e)}")

            # --- 设置文档和样式 ---
            doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            style_n = styles['Normal']
            style_n.fontName = font_name  # 应用中文字体
            style_n.fontSize = 12
            style_n.leading = 16  # 行间距
            
            story = []
            # 将传入的字符串路径拆分为列表
            path_list = [p.strip() for p in image_paths.split(',') if p.strip()]

            if render_markdown:
                story.extend(self._markdown_to_pdf_story(content, styles, font_name, path_list))
            else:
                # --- 处理文本与图片占位符 ---
                parts = re.split(r'(\{\{image_\d+\}\})', content)

                for part in parts:
                    placeholder_match = re.match(r'\{\{image_(\d+)\}\}', part)
                    
                    if placeholder_match:
                        idx = int(placeholder_match.group(1)) - 1
                        if 0 <= idx < len(path_list):
                            img_path = path_list[idx]
                            if os.path.exists(img_path):
                                img = RLImage(img_path)
                                max_width = 5 * inch
                                if img.drawWidth > max_width:
                                    ratio = max_width / img.drawWidth
                                    img.drawWidth = max_width
                                    img.drawHeight = img.drawHeight * ratio
                                story.append(Spacer(1, 12))
                                story.append(img)
                                story.append(Spacer(1, 12))
                            else:
                                story.append(Paragraph(f"[图片文件未找到: {img_path}]", style_n))
                        else:
                            story.append(Paragraph(f"[缺失对应的图片路径: {part}]", style_n))
                    else:
                        if part.strip():
                            formatted_text = part.replace('\n', '<br/>')
                            story.append(Paragraph(formatted_text, style_n))

            # --- 生成 PDF ---
            doc.build(story)
            return File(output_path)

        except Exception as e:
            return f'An error occurred: {str(e)}'

    def _markdown_to_pdf_story(self, md_text: str, styles, font_name: str, path_list: List[str]):
        """将 Markdown 文本解析为 reportlab Flowables 列表。仅使用 Paragraph/List/Preformatted 等基础组件。"""
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            Paragraph,
            Image as RLImage,
            Spacer,
            ListFlowable,
            ListItem,
            Preformatted,
        )
        story: List = []
        normal = styles['Normal']
        normal.fontName = font_name

        heading_styles = {
            1: styles.get('Heading1', normal.clone('Heading1')),
            2: styles.get('Heading2', normal.clone('Heading2')),
            3: styles.get('Heading3', normal.clone('Heading3')),
            4: styles.get('Heading4', normal.clone('Heading4')),
            5: styles.get('Heading5', normal.clone('Heading5')),
            6: styles.get('Heading6', normal.clone('Heading6')),
        }
        for s in heading_styles.values():
            try:
                s.fontName = font_name
            except Exception:
                pass

        lines = md_text.splitlines()
        i = 0

        def inline_html(txt: str) -> str:
            # 链接 -> 文本 (url)
            txt = re.sub(r"\[([^\]]+)\]\(([^\)]+)\)", lambda m: f"{m.group(1)} ({m.group(2)})", txt)
            # 粗体、斜体、行内代码 -> 受限 HTML
            txt = re.sub(r"\*\*([^*]+)\*\*", r"<b>\1</b>", txt)
            txt = re.sub(r"\*([^*]+)\*", r"<i>\1</i>", txt)
            txt = re.sub(r"`([^`]+)`", r"<font face='Courier'>\1</font>", txt)
            # 换行
            return txt

        while i < len(lines):
            line = lines[i]
            if not line.strip():
                i += 1
                continue

            # 代码块
            if line.strip().startswith('```'):
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith('```'):
                    code_lines.append(lines[i])
                    i += 1
                i += 1
                story.append(Preformatted('\n'.join(code_lines), normal))
                story.append(Spacer(1, 8))
                continue

            # 标题
            m = re.match(r"^(#{1,6})\s+(.*)$", line)
            if m:
                level = min(len(m.group(1)), 6)
                txt = inline_html(m.group(2))
                story.append(Paragraph(txt, heading_styles[level]))
                story.append(Spacer(1, 6))
                i += 1
                continue

            # 列表
            if re.match(r"^\s*([*+-])\s+", line) or re.match(r"^\s*\d+\.\s+", line):
                is_ordered = bool(re.match(r"^\s*\d+\.\s+", line))
                items = []
                while i < len(lines) and (re.match(r"^\s*([*+-])\s+", lines[i]) or re.match(r"^\s*\d+\.\s+", lines[i])):
                    item_text = re.sub(r"^\s*([*+-]|\d+\.)\s+", "", lines[i]).strip()
                    items.append(ListItem(Paragraph(inline_html(item_text), normal)))
                    i += 1
                story.append(ListFlowable(items, bulletType='1' if is_ordered else 'bullet'))
                story.append(Spacer(1, 6))
                continue

            # 图片行：![](path)
            img_md = re.match(r"^!\[[^\]]*\]\(([^\)]+)\)", line.strip())
            ph_md = re.match(r"^\{\{image_(\d+)\}\}$", line.strip())
            if img_md:
                img_path = img_md.group(1)
                if os.path.exists(img_path):
                    img = RLImage(img_path)
                    max_width = 5 * inch
                    if img.drawWidth > max_width:
                        ratio = max_width / img.drawWidth
                        img.drawWidth = max_width
                        img.drawHeight = img.drawHeight * ratio
                    story.append(img)
                    story.append(Spacer(1, 12))
                else:
                    story.append(Paragraph(f"[图片文件未找到: {img_path}]", normal))
                i += 1
                continue
            if ph_md:
                idx = int(ph_md.group(1)) - 1
                if 0 <= idx < len(path_list):
                    img_path = path_list[idx]
                    if os.path.exists(img_path):
                        img = RLImage(img_path)
                        max_width = 5 * inch
                        if img.drawWidth > max_width:
                            ratio = max_width / img.drawWidth
                            img.drawWidth = max_width
                            img.drawHeight = img.drawHeight * ratio
                        story.append(img)
                        story.append(Spacer(1, 12))
                    else:
                        story.append(Paragraph(f"[图片文件未找到: {img_path}]", normal))
                else:
                    story.append(Paragraph(f"[缺失对应的图片路径: {{image_{idx+1}}}]", normal))
                i += 1
                continue

            # 普通段落，支持行内占位符
            parts = re.split(r'(\{\{image_\d+\}\})', line)
            if len(parts) == 1:
                story.append(Paragraph(inline_html(line), normal))
                story.append(Spacer(1, 4))
            else:
                # 先把文本段落放进去，再插图片，再文本...
                buf = []
                for part in parts:
                    ph = re.match(r'\{\{image_(\d+)\}\}', part)
                    if ph:
                        if buf:
                            story.append(Paragraph(inline_html(''.join(buf)), normal))
                            story.append(Spacer(1, 4))
                            buf = []
                        idx = int(ph.group(1)) - 1
                        if 0 <= idx < len(path_list):
                            img_path = path_list[idx]
                            if os.path.exists(img_path):
                                img = RLImage(img_path)
                                max_width = 5 * inch
                                if img.drawWidth > max_width:
                                    ratio = max_width / img.drawWidth
                                    img.drawWidth = max_width
                                    img.drawHeight = img.drawHeight * ratio
                                story.append(img)
                                story.append(Spacer(1, 12))
                            else:
                                story.append(Paragraph(f"[图片文件未找到: {img_path}]", normal))
                        else:
                            story.append(Paragraph(f"[缺失对应的图片路径: {part}]", normal))
                    else:
                        buf.append(part)
                if buf:
                    story.append(Paragraph(inline_html(''.join(buf)), normal))
                    story.append(Spacer(1, 4))

            # IMPORTANT: advance to next line to avoid infinite loop.
            i += 1

        return story


# class PdfFileGenerator(BaseTool):
#     default_desc = 'Generates a PDF file with content, support EN / ZH font and optional image insertion.'

#     def apply(self, filename: str, content: str, image_path: str = None, output_dir: str = None) -> Annotated[File, Info('Generated PDF file.', filetype='application/pdf')]:
#         try:
#             if output_dir is None:
#                 output_dir = os.path.join(os.getcwd(), "tools_output")
#             os.makedirs(output_dir, exist_ok=True)
#             output_path = os.path.join(output_dir, filename + ".pdf")

#             # 1. 注册中文字体 (请确保路径正确)
#             # 这里以常见的 SimSun (宋体) 为例，你可以根据实际情况指定 .ttf 文件位置
#             font_path = "simsun.ttc"  # 请确保该文件在当前目录或指定绝对路径
            
#             # 只有在还没注册时才注册，避免重复注册报错
#             try:
#                 pdfmetrics.registerFont(TTFont('SimSun', font_path))
#                 font_name = 'SimSun'
#             except:
#                 # 如果没找到指定字体，回退到系统默认（注意：默认 Helvetica 不支持中文，会显示为方块）
#                 font_name = 'Helvetica'
#                 print(f"Warning: Font file not found at {font_path}, fallback to Helvetica.")

#             # 2. 初始化文档和样式
#             doc = SimpleDocTemplate(output_path, pagesize=letter)
#             styles = getSampleStyleSheet()
#             style_n = styles['Normal']
#             style_n.fontName = font_name
            
#             # 用于存放文档元素的列表
#             story = []

#             # 3. 添加图片 (如果提供了路径)
#             if image_path and os.path.exists(image_path):
#                 # 创建图片对象
#                 # width 和 height 可以手动指定，也可以自适应
#                 # 这里设置宽度为 4 英寸，高度根据比例自动缩放
#                 img = Image(image_path)
                
#                 # 限制图片最大宽度，防止溢出页面
#                 max_width = 6 * inch
#                 if img.drawWidth > max_width:
#                     ratio = max_width / img.drawWidth
#                     img.drawWidth = max_width
#                     img.drawHeight = img.drawHeight * ratio
                
#                 story.append(img)
#                 story.append(Spacer(1, 12))  # 图片和文字之间加一点间距 (12pt)

#             # 4. 添加文本内容
#             formatted_content = content.replace('\n', '<br/>')
#             story.append(Paragraph(formatted_content, style_n))

#             # 5. 构建 PDF
#             doc.build(story)

#             return File(output_path, filetype='application/pdf')
#         except Exception as e:
#             return f'An error occurred: {str(e)}'

# class PdfFileGenerator(BaseTool):
#     default_desc = 'Generates a PDF file with the provided content.'

#     def apply(self, filename: str, content: str, output_dir: str = None) -> Annotated[File, Info('Generated PDF file.', filetype='application/pdf')]:
#         import os
#         try:
#             if output_dir is None:
#                 output_dir = os.path.join(os.getcwd(), "tools_output")
#             os.makedirs(output_dir, exist_ok=True)
#             output_path = os.path.join(output_dir, filename + ".pdf")
#             c = canvas.Canvas(output_path, pagesize=letter)
#             width, height = letter
#             c.drawString(100, height - 100, content)
#             c.save()
#             return File(output_path, filetype='application/pdf')
#         except Exception as e:
#             return f'An error occurred: {str(e)}'



class PptxFileGenerator(BaseTool):
    default_desc = (
        'Generates a PowerPoint (.pptx) presentation supporting Chinese characters and image embedding. '
        'The "slides_content_str" should be a JSON string of a list of objects, each with "title" and "content". '
        'In "content", use {{image_1}}, {{image_2}} placeholders to insert images. '
        'The "image_paths" parameter must be a comma-separated string of file paths matching the placeholder order.'
    )

    def apply(
        self, 
        filename: str, 
        slides_content_str: str, 
        image_paths: str = "", 
        output_dir: str = None
    ) -> Annotated[File, Info('Generated PPTX file.', filetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')]:
        from pptx import Presentation
        from pptx.oxml.ns import qn
        from pptx.util import Inches
        
        if output_dir is None:
            output_dir = DEFAULT_TOOLS_OUTPUT_DIR
        os.makedirs(output_dir, exist_ok=True)
        
        path_list = [p.strip() for p in image_paths.split(',') if p.strip()]

        try:
            slides_data = json.loads(slides_content_str)
        except json.JSONDecodeError:
            return 'Invalid format: slides_content_str must be a JSON list of objects.'

        try:
            output_path = os.path.join(output_dir, filename + ".pptx")
            prs = Presentation()
            cn_font = 'SimSun'

            for slide_item in slides_data:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                # 1. 处理标题
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = slide_item.get("title", "")
                    for paragraph in title_shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            self._set_font(run, cn_font)

                # 2. 处理内容区域
                content_body = slide_item.get("content", "")
                placeholder = slide.placeholders[1]
                placeholder.text = "" 
                
                parts = re.split(r'(\{\{image_\d+\}\})', content_body)
                tf = placeholder.text_frame
                tf.word_wrap = True

                for part in parts:
                    placeholder_match = re.match(r'\{\{image_(\d+)\}\}', part)
                    if placeholder_match:
                        idx = int(placeholder_match.group(1)) - 1
                        if 0 <= idx < len(path_list):
                            img_path = path_list[idx]
                            if os.path.exists(img_path):
                                # 默认图片位置
                                slide.shapes.add_picture(img_path, Inches(1), Inches(3.5), width=Inches(4))
                    else:
                        if part.strip():
                            p = tf.add_paragraph()
                            p.text = part.strip()
                            for run in p.runs:
                                self._set_font(run, cn_font)

            prs.save(output_path)
            return File(output_path)
            
        except Exception as e:
            import traceback
            return f'An error occurred: {str(e)}\n{traceback.format_exc()}'

    def _set_font(self, run, font_name):
        """稳健地设置中文字体属性"""
        from pptx.oxml.ns import qn
        run.font.name = font_name
        # 获取或创建 rPr 节点
        # 使用 hasattr 检查以避免 CT_TextCharacterProperties 报错
        el = run.font._element
        rPr = el.get_or_add_rPr() if hasattr(el, 'get_or_add_rPr') else None
        
        if rPr is not None:
            # 设置东亚字体
            ea = rPr.get_or_add_ea() if hasattr(rPr, 'get_or_add_ea') else None
            # 直接通过 XML 属性设置，绕过部分接口缺失问题
            rPr.set(qn('w:eastAsia'), font_name)


# class PptxFileGenerator(BaseTool):
#     default_desc = 'Generates a PowerPoint presentation with the provided slides content.'

#     def apply(self, filename: str, slides_content_str: str, output_dir: str = None) -> Annotated[File, Info('Generated PPTX file.', filetype='office/pptx')]:
#         import os
#         from pptx import Presentation

#         if output_dir is None:
#             output_dir = os.path.join(os.getcwd(), "tools_output")
#         os.makedirs(output_dir, exist_ok=True)

#         try:
#             slides = json.loads(slides_content_str)
#         except json.JSONDecodeError:
#             return 'Invalid slides content format: data should be a JSON-formatted string representing a list of slides.'

#         try:
#             output_path = os.path.join(output_dir, filename + ".pptx")
#             prs = Presentation()
#             for slide in slides:
#                 s = prs.slides.add_slide(prs.slide_layouts[1])
#                 # set title if available
#                 try:
#                     s.shapes.title.text = slide.get("title", "")
#                 except Exception:
#                     pass
#                 # put main content into the first placeholder if present
#                 try:
#                     s.placeholders[1].text = slide.get("content", "")
#                 except Exception:
#                     # ignore if the placeholder isn't available
#                     pass
#             prs.save(output_path)
#             return File(output_path, filetype='office/pptx')
#         except Exception as e:
#             return f'An error occurred: {str(e)}'


class XlsxFileGenerator(BaseTool):
    default_desc = 'Generates an XLSX file with the provided data.'

    def apply(self, filename: str, data_str: str, output_dir: str = None) -> Annotated[File, Info('Generated XLSX file.', filetype='office/xlsx')]:
        import os
        try:
            data = json.loads(data_str)
            if output_dir is None:
                output_dir = DEFAULT_TOOLS_OUTPUT_DIR
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, filename + ".xlsx")
            wb = openpyxl.Workbook()
            ws = wb.active
            for row in data:
                ws.append(row)
            wb.save(output_path)
            return File(output_path, filetype='office/xlsx')
        except json.JSONDecodeError:
            return 'Invalid data format: data should be a JSON-formatted string representing a list of lists.'
        except Exception as e:
            return f'An error occurred: {str(e)}'


class HtmlFileGenerator(BaseTool):
    default_desc = (
        'Generates an HTML file from provided markup. Supports fenced code blocks '
        'like ```html ... ``` or raw HTML. If no <html> tag is present, the tool '
        'wraps the content into a minimal HTML5 skeleton with UTF-8 meta.'
    )

    def apply(
        self,
        filename: str,
        html_code: str,
        output_dir: str = None,
    ) -> Annotated[File, Info('Generated HTML file.', filetype='text/html')]:
        try:
            # Extract code from fenced blocks if provided
            if '```html' in html_code:
                html_code = html_code.split('```html')[1].split('```')[0]
            elif '```' in html_code:
                html_code = html_code.split('```')[1].split('```')[0]

            content = html_code.strip()

            # Wrap with basic HTML skeleton if missing root tag
            if '<html' not in content.lower():
                content = (
                    '<!doctype html>\n'
                    '<html lang="en">\n'
                    '  <head>\n'
                    '    <meta charset="utf-8" />\n'
                    '    <meta name="viewport" content="width=device-width, initial-scale=1" />\n'
                    '    <title>Document</title>\n'
                    '    <style>body{font-family:system-ui,-apple-system,Segoe UI,Roboto,Helvetica,Arial;line-height:1.5;margin:2rem;}</style>\n'
                    '  </head>\n'
                    '  <body>\n'
                    f'{content}\n'
                    '  </body>\n'
                    '</html>\n'
                )

            # Determine output directory
            if output_dir is None:
                output_dir = DEFAULT_TOOLS_OUTPUT_DIR
            os.makedirs(output_dir, exist_ok=True)

            # Ensure .html suffix
            out_name = filename if filename.endswith('.html') else filename + '.html'
            output_path = os.path.join(output_dir, out_name)

            # Write file as UTF-8
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(content)

            return File(output_path, filetype='text')
        except Exception as e:
            return f'An error occurred: {str(e)}'


class ImageQualityTest(BaseTool):
    """
    A tool to evaluate image quality by measuring the Laplacian variance of the input image.
    
    Args:
        threshold (float): The variance threshold below which the image is considered blurry.
            Defaults to 100.0.
        toolmeta (None | dict | ToolMeta): Additional tool metadata. Defaults to None.
    """

    default_desc = 'This tool evaluates image quality by calculating the Laplacian variance, which is a proxy for image sharpness.'

    @require('opencv-python')
    def __init__(self, threshold: float = 100.0, toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.threshold = threshold

    def setup(self):
        # OpenCV 和 numpy 均已安装，无需额外初始化操作。
        pass

    def apply(
        self,
        image: ImageIO,
    ) -> Annotated[str, Info('Image quality test result including Laplacian variance and a flag indicating if the image is blurry.')]:
        # 将输入转换为 numpy 数组
        img_array = image.to_array()

        # 如果是彩色图像，则转换为灰度图
        if len(img_array.shape) == 3:
            gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        else:
            gray = img_array

        # 计算 Laplacian 以评估图像的锐度
        laplacian = cv2.Laplacian(gray, cv2.CV_64F)
        variance = laplacian.var()

        # 判断图像是否模糊
        is_blurry = variance < self.threshold

        # 格式化输出结果
        output = "Laplacian variance: {:.2f}, Blurry: {}".format(variance, is_blurry)
        return output


class NoiseReduction(BaseTool):
    """
    A tool for reducing noise in audio files using the 'noisereduce' library.

    Args:
        toolmeta (None | dict | ToolMeta): Additional information about the tool, default is None.
    """

    default_desc = "This tool reduces noise in input audio files using spectral gating."

    def __init__(self, toolmeta=None):
        super().__init__(toolmeta=toolmeta)

    def apply(
        self,
        audio: str,
        output_dir: str = None
    ) -> Annotated[AudioIO, Info("Audio file with reduced noise.")]:
        # Input is expected to be a file path string.
        audio_path = audio.to_path() if hasattr(audio, 'to_path') else str(audio)
        data, rate = sf.read(audio_path)

        # Apply noise reduction
        reduced_noise = nr.reduce_noise(y=data, sr=rate)

        # Save the denoised audio to a new file
        import os
        if output_dir is None:
            output_dir = DEFAULT_TOOLS_OUTPUT_DIR
        os.makedirs(output_dir, exist_ok=True)
        output_file = "denoised.wav"
        output_path = os.path.join(output_dir, output_file)
        sf.write(output_path, reduced_noise, rate)
        return AudioIO(output_path)


class PitchShifting(BaseTool):
    """
    A tool to perform pitch shifting on audio files.
    
    Args:
        n_steps (float): Number of semitones to shift the pitch.
                         Positive values shift the pitch up, negative values shift down.
                         Defaults to 0.0.
        sample_rate (int): The target sample rate for processing. Defaults to 22050.
        toolmeta (None | dict | ToolMeta): Additional tool metadata. Defaults to None.
    """

    default_desc = 'This tool applies pitch shifting to an audio file, altering its pitch by the specified number of semitones.'

    @require('librosa')
    def __init__(self, n_steps: float = 0.0, sample_rate: int = 22050, toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.n_steps = n_steps
        self.sample_rate = sample_rate

    def setup(self):
        # 确保 librosa 已加载
        import librosa
        self._librosa = librosa

    def apply(
        self,
        audio: str,
        output_dir: str = None
    ) -> Annotated[AudioIO, Info('Audio file with pitch shifted by n_steps semitones.')]:
        """
        Apply pitch shifting to the input audio.
        
        Args:
            audio (AudioIO): The input audio object.
        
        Returns:
            AudioIO: The audio object after pitch shifting.
        """
        # Input is expected to be a file path string.
        audio_path = audio.to_path() if hasattr(audio, 'to_path') else str(audio)
        y, sr = sf.read(audio_path)

        # 如果采样率不匹配，则重采样
        if sr != self.sample_rate:
            y = self._librosa.resample(y, orig_sr=sr, target_sr=self.sample_rate)
            sr = self.sample_rate

        # 使用关键字参数调用 pitch_shift
        y_shifted = self._librosa.effects.pitch_shift(y=y, sr=sr, n_steps=self.n_steps)

        # 构造输出文件路径，例如将 "sample.wav" 转为 "sample_pitchshift_+2.0.wav"
        base_path, ext = audio_path.rsplit('.', 1)
        import os
        if output_dir is None:
            output_dir = DEFAULT_TOOLS_OUTPUT_DIR
        os.makedirs(output_dir, exist_ok=True)
        output_filename = f"{os.path.basename(base_path)}_pitchshift_{self.n_steps:+.1f}.{ext}"
        output_path = os.path.join(output_dir, output_filename)
        sf.write(output_path, y_shifted, sr)
        return AudioIO(output_path)



class SoundCharacteristicExtraction(BaseTool):
    """
    A tool to extract sound characteristics from an input audio file.
    
    This tool uses librosa to extract the following features:
      - MFCC (Mel Frequency Cepstral Coefficients)
      - Spectral Centroid
      - Spectral Rolloff
      - Zero Crossing Rate

    For each feature, the tool computes the mean and variance along the time axis.
    
    Args:
        sample_rate (int): The target sample rate for processing. Defaults to 22050.
        n_mfcc (int): Number of MFCC coefficients to extract. Defaults to 13.
        toolmeta (None | dict | ToolMeta): Additional metadata for the tool. Defaults to None.
    """

    default_desc = (
        "This tool extracts sound characteristics (MFCC, spectral centroid, spectral rolloff, "
        "and zero crossing rate) from an audio file and returns the mean and variance of each feature in JSON format."
    )

    @require('librosa')
    def __init__(self, sample_rate: int = 22050, n_mfcc: int = 13, toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.sample_rate = sample_rate
        self.n_mfcc = n_mfcc

    def setup(self):
        """
        Initialize the tool by importing librosa and storing it for later use.
        """
        import librosa
        self._librosa = librosa

    def apply(
        self,
        audio: str,
    ) -> Annotated[str, Info('Extracted sound characteristics as a JSON string.')]:
        """
        Extract sound characteristics from the input audio.
        
        Args:
            audio (AudioIO): The input audio object.
        
        Returns:
            str: A JSON string containing the mean and variance of extracted features:
                 - MFCC (list of means and variances for each coefficient)
                 - Spectral Centroid
                 - Spectral Rolloff
                 - Zero Crossing Rate
        """
        # Input is expected to be a file path string.
        audio_path = audio.to_path() if hasattr(audio, 'to_path') else str(audio)
        y, sr = self._librosa.load(audio_path, sr=self.sample_rate)

        # Compute MFCC features (shape: (n_mfcc, frames))
        mfcc = self._librosa.feature.mfcc(y=y, sr=sr, n_mfcc=self.n_mfcc)
        mfcc_mean = mfcc.mean(axis=1).tolist()
        mfcc_var = mfcc.var(axis=1).tolist()

        # Compute spectral centroid (shape: (1, frames))
        spectral_centroid = self._librosa.feature.spectral_centroid(y=y, sr=sr)
        spectral_centroid_mean = float(spectral_centroid.mean())
        spectral_centroid_var = float(spectral_centroid.var())

        # Compute spectral rolloff (shape: (1, frames))
        spectral_rolloff = self._librosa.feature.spectral_rolloff(y=y, sr=sr)
        spectral_rolloff_mean = float(spectral_rolloff.mean())
        spectral_rolloff_var = float(spectral_rolloff.var())

        # Compute zero crossing rate (shape: (1, frames))
        zero_crossing_rate = self._librosa.feature.zero_crossing_rate(y)
        zero_crossing_rate_mean = float(zero_crossing_rate.mean())
        zero_crossing_rate_var = float(zero_crossing_rate.var())

        # Construct the features dictionary
        features = {
            "mfcc": {
                "mean": mfcc_mean,
                "var": mfcc_var
            },
            "spectral_centroid": {
                "mean": spectral_centroid_mean,
                "var": spectral_centroid_var
            },
            "spectral_rolloff": {
                "mean": spectral_rolloff_mean,
                "var": spectral_rolloff_var
            },
            "zero_crossing_rate": {
                "mean": zero_crossing_rate_mean,
                "var": zero_crossing_rate_var
            }
        }

        # Return the JSON string of extracted features
        return json.dumps(features, indent=2)


class DiaSpeakerrization(BaseTool):
    """A tool to perform speaker diarization on an audio file.

    Args:
        model (str): The pre-trained model identifier for speaker diarization.
            Defaults to 'pyannote/speaker-diarization'.
        toolmeta (None | dict | ToolMeta): Additional info for the tool.
            Defaults to None.
        **model_args: Additional keyword arguments for model initialization.
            If a 'pipeline_name' parameter is provided, it will override the model.
            To use a gated model, provide a valid Hugging Face token via `use_auth_token`.
    """

    default_desc = 'This tool performs speaker diarization on the input audio file.'

    @require('pyannote.audio')
    def __init__(self,
                 model: str = 'pyannote/speaker-diarization',
                 toolmeta=None,
                 **model_args):
        super().__init__(toolmeta=toolmeta)
        # 如果传入了 pipeline_name，则用其作为 model 参数
        if 'pipeline_name' in model_args:
            model = model_args.pop('pipeline_name')
        self.model = model
        self.model_args = model_args

    def setup(self):
        from pyannote.audio import Pipeline
        try:
            # 尝试加载预训练模型，同时传入额外参数（例如 use_auth_token）
            self._pipeline = Pipeline.from_pretrained(self.model, **self.model_args)
            if self._pipeline is None:
                raise RuntimeError("Failed to load pyannote.audio Pipeline.")
        except Exception as e:
            raise RuntimeError(f"Error loading pyannote.audio Pipeline: {e}")

    def apply(
        self,
        audio: str,
    ) -> Annotated[str,
                   Info("Diarization results with each segment's start and end times and the corresponding speaker label.")]:
        # Input is expected to be a file path string.
        audio_path = audio.to_path() if hasattr(audio, 'to_path') else str(audio)
        diarization = self._pipeline(audio_path)
        outputs = []
        # 遍历分离结果，构造结果字符串
        for segment, _, speaker in diarization.itertracks(yield_label=True):
            outputs.append("({:.2f}, {:.2f}) Speaker: {}".format(segment.start, segment.end, speaker))
        return "\n".join(outputs)


class ReadCSV(BaseTool):
    """A tool to read CSV files and return their content as a CSV string.

    Args:
        delimiter (str): The delimiter used in the CSV file. Defaults to ','.
        toolmeta (Optional[dict]): Additional tool metadata. Defaults to None.
    """

    default_desc = "This tool reads a CSV file and returns its content as a CSV string."

    def __init__(self, delimiter: str = ',', toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.delimiter = delimiter

    def apply(
        self,
        csv_file: File,
    ) -> Annotated[str, Info("CSV content as a string.")]:
        # Retrieve the file path using the File IO utility.
        csv_path = csv_file.to_path()
        # Read the CSV file using pandas with the specified delimiter.
        df = pd.read_csv(csv_path, delimiter=self.delimiter)
        # Convert the DataFrame to CSV text (without the index)
        csv_text = df.to_csv(index=False)
        return csv_text

class ReadDOCX(BaseTool):
    """A tool to read DOCX files and return their text content.

    Args:
        toolmeta (Optional[dict]): Additional tool metadata. Defaults to None.
    """
    default_desc = "This tool reads a DOCX file and returns its text content."

    def apply(
        self,
        file: File,
    ) -> Annotated[str, Info("Extracted text content from DOCX file.")]:
        # Retrieve the file path from the File IO object.
        docx_path = file.to_path()
        import docx
        # Open the DOCX document.
        doc = docx.Document(docx_path)
        # Extract text from each paragraph.
        full_text = "\n".join([para.text for para in doc.paragraphs])
        return full_text


class ReadPDF(BaseTool):
    """A tool to read PDF files and return their text content.

    Args:
        toolmeta (Optional[dict]): Additional tool metadata.
            Defaults to None.
    """
    default_desc = "This tool reads a PDF file and returns its text content."

    def apply(
        self,
        file: File,
    ) -> Annotated[str, Info("Extracted text from the PDF file.")]:
        # Get file path from File object.
        pdf_path = file.to_path()
        reader = PyPDF2.PdfReader(pdf_path)
        full_text = []
        for page in reader.pages:
            full_text.append(page.extract_text())
        return "\n".join(full_text)

class ReadPPTX(BaseTool):
    """A tool to read PPTX files and return their text content.

    Args:
        toolmeta (Optional[dict]): Additional tool metadata.
            Defaults to None.
    """
    default_desc = "This tool reads a PPTX file and returns its text content."

    def apply(
        self,
        file: File,
    ) -> Annotated[str, Info("Extracted text from the PPTX file.")]:
        # Get file path from File object.
        pptx_path = file.to_path()
        from pptx import Presentation
        prs = Presentation(pptx_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)

class ReadXLSX(BaseTool):
    """A tool to read XLSX files and return their content as CSV-formatted text.

    Args:
        toolmeta (Optional[dict]): Additional tool metadata.
            Defaults to None.
    """
    default_desc = "This tool reads an XLSX file and returns its content as CSV text."

    def __init__(self, delimiter: str = ",", toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.delimiter = delimiter

    def apply(
        self,
        file: File,
    ) -> Annotated[str, Info("Extracted XLSX content as CSV text.")]:
        # Get file path from File object.
        xlsx_path = file.to_path()
        # Read the Excel file using pandas.
        df = pd.read_excel(xlsx_path)
        # Convert DataFrame to CSV-formatted string (without index)
        csv_text = df.to_csv(index=False, sep=self.delimiter)
        return csv_text


class Prover(BaseTool):
    """
    一个辅助数学定理证明的工具。

    该工具要求输入的定理为 SMT-LIB2 格式的逻辑公式字符串。
    工具通过将定理 F 的否定（Not(F)）加入 Z3 求解器，
    若其不可满足（unsat），则证明 F 为重言式，从而定理成立；否则返回一个反例。

    示例：证明 "forall x:Int, x + 1 > x" 可以使用如下 SMT-LIB2 格式：
    
    (declare-const x Int)
    (assert (forall ((x Int)) (> (+ x 1) x)))
    """
    
    default_desc = "该工具利用 Z3 SMT 求解器辅助证明数学定理（输入为 SMT-LIB2 格式字符串）。"

    def __init__(self, toolmeta=None):
        super().__init__(toolmeta=toolmeta)

    def setup(self):
        import z3
        self.z3 = z3

    def apply(
        self,
        theorem: str,
    ) -> Annotated[str, Info("证明结果：若定理成立则返回证明成功；否则返回反例。")]:
        """
        尝试证明输入的定理是否成立。

        参数:
            theorem (str): SMT-LIB2 格式的逻辑公式字符串。

        返回:
            str: 如果证明成功，返回 "证明成功：定理为重言式。"；
                 如果有反例，则返回 "定理不成立，反例：..."。
        """
        z3 = self.z3
        try:
            # 解析 SMT-LIB2 格式字符串为一组断言
            assertions = z3.parse_smt2_string(theorem)
        except Exception as e:
            return f"解析定理失败: {e}"
        
        # 将所有断言合取为公式 F
        F = z3.And(assertions)
        solver = z3.Solver()
        # 将 Not(F) 加入求解器，用于判断 F 是否为重言式
        solver.add(z3.Not(F))
        result = solver.check()
        if result == z3.unsat:
            return "证明成功：定理为重言式。"
        elif result == z3.sat:
            return f"定理不成立，反例：{solver.model()}"
        else:
            return "证明结果未知。"



class VideoClip(BaseTool):
    """A tool to extract a clip from a video between specified start and end times.

    Args:
        start_time (float): The start time of the clip in seconds.
        end_time (float): The end time of the clip in seconds.
        output_format (str): The format of the output video (e.g., 'mp4').
        device (str | bool): The device to load the model. Defaults to True, meaning automatic selection.
        toolmeta (None | dict | ToolMeta): Additional tool metadata. Defaults to None.
    """

    default_desc = 'This tool extracts a clip from a video between specified start and end times.'

    @require('opencv-python')
    def __init__(self,
                 start_time: float = 0.0,
                 end_time: float = 10.0,
                 output_format: str = 'mp4',
                 device: Union[bool, str] = True,
                 toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.start_time = start_time
        self.end_time = end_time
        self.output_format = output_format
        self.device = device

    def setup(self):
        # Perform any necessary setup here
        pass

    def apply(self,
              video: str,
              start_time: Optional[float] = None,
              end_time: Optional[float] = None,
              filename: Optional[str] = None,
              output_dir: str = None
              ) -> VideoIO:
        """
        Extract a clip from `video`. start_time and end_time (seconds) may be
        passed here; if omitted the constructor defaults are used.
        """
        # choose call-time values if provided
        st = self.start_time if start_time is None else start_time
        ed = self.end_time if end_time is None else end_time

        # validate
        if st is None or ed is None:
            raise ValueError("start_time and end_time must be provided either in constructor or apply() call")
        if st < 0 or ed <= st:
            raise ValueError("Invalid start_time/end_time: must be non-negative and end_time > start_time")

        # Input is expected to be a file path string.
        video_path = video.to_path() if hasattr(video, 'to_path') else str(video)

        cap = cv2.VideoCapture(video_path)
        if not cap.isOpened():
            raise ValueError(f"Cannot open video file: {video_path}")

        fps = cap.get(cv2.CAP_PROP_FPS)
        if fps is None or fps <= 0:
            cap.release()
            raise RuntimeError("Failed to get FPS from video; cannot compute frames for clipping")

        frame_width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        frame_height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT)) if cap.get(cv2.CAP_PROP_FRAME_COUNT) > 0 else None

        # Compute start and end frames (frame indices)
        start_frame = int(st * fps)
        end_frame = int(ed * fps) - 1  # treat end_time as exclusive; clip up to the previous frame
        if total_frames is not None:
            end_frame = min(end_frame, total_frames - 1)

        # Prepare output path
        import os
        if output_dir is None:
            output_dir = DEFAULT_TOOLS_OUTPUT_DIR
        os.makedirs(output_dir, exist_ok=True)

        # normalize extension and build filename
        ext = str(self.output_format).lstrip('.')
        if filename is None:
            safe_base = f"clip_{st:.3f}_{ed:.3f}".replace('.', '_')
            output_filename = f"{safe_base}.{ext}"
        else:
            # ensure only basename is used and extension is correct
            filename = os.path.basename(filename)
            if not filename.lower().endswith(f".{ext}"):
                output_filename = f"{filename}.{ext}"
            else:
                output_filename = filename

        output_path = os.path.join(output_dir, output_filename)

        fourcc = cv2.VideoWriter_fourcc(*'mp4v')
        out = cv2.VideoWriter(output_path, fourcc, fps, (frame_width, frame_height))

        # Seek to start frame
        cap.set(cv2.CAP_PROP_POS_FRAMES, start_frame)
        current_frame = start_frame
        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break
            if current_frame > end_frame:
                break
            out.write(frame)
            current_frame += 1

        cap.release()
        out.release()
        return VideoIO(output_path)


class VideoDescription(BaseTool):
    """A tool to generate descriptions for a video by captioning key frames.

    Args:
        model_name (str): The image-to-text model to use for captioning.
            Defaults to 'Salesforce/blip-image-captioning-large'.
        frame_skip (int): Process every Nth frame for efficiency.
            Defaults to 5.
        toolmeta (Optional[dict]): Additional tool metadata. Defaults to None.
    """

    default_desc = 'This tool generates video descriptions by captioning key frames.'

    @require('transformers')
    def __init__(self, model_name: str = 'Salesforce/blip-image-captioning-large',
                 frame_skip: int = 5,
                 toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.model_name = model_name
        self.frame_skip = frame_skip

    def setup(self):
        from transformers import pipeline
        # Load the image captioning pipeline from Hugging Face.
        self.captioner = load_or_build_object(pipeline, "image-to-text", model=self.model_name)

    def apply(
        self,
        video: str,
        frame_skip: Optional[int] = None,
    ) -> Annotated[str, Info('Generated video description by captioning key frames.')]:
        """
        Generate captions for key frames. `frame_skip` may be provided here to
        override the constructor default (must be >= 1).
        """
        fs = self.frame_skip if frame_skip is None else frame_skip

        video_path = video.to_path() if hasattr(video, 'to_path') else str(video)
        cap = cv2.VideoCapture(video_path)
        frame_count = 0
        captions = []

        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break

            if frame_count % fs == 0:
                # Convert the frame from BGR to RGB and then to a PIL Image
                frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                pil_image = PILImage.fromarray(frame_rgb)

                # Generate caption using the image-to-text pipeline
                caption_result = self.captioner(pil_image)
                caption = caption_result[0]["generated_text"] if caption_result else ""
                captions.append(f"Frame {frame_count}: {caption}")

            frame_count += 1

        cap.release()
        # Optionally, you can aggregate or summarize the captions further.
        return "\n".join(captions)


class VideoObjectDetection(BaseTool):
    """A tool to perform object detection on a video frame by frame.

    Args:
        model_path (str): The path or identifier for the YOLO model.
            Defaults to 'yolov8n.pt'.
        frame_skip (int): Process every Nth frame for efficiency.
            Defaults to 1 (process all frames).
        toolmeta (Optional[dict]): Additional tool metadata.
            Defaults to None.
    """

    default_desc = 'This tool detects objects in video frames using a YOLO model.'

    @require('ultralytics')
    def __init__(self,
                 model_path: str = 'yolov8n.pt',
                 frame_skip: int = 1,
                 toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.model_path = model_path
        self.frame_skip = frame_skip

    def setup(self):
        from ultralytics import YOLO
        # Load the YOLO object detection model using the ultralytics YOLO API.
        self.model = load_or_build_object(YOLO, self.model_path)

    def apply(
        self,
        video: str,
        frame_skip: Optional[int] = None,
    ) -> Annotated[str, Info('Object detection results for each processed video frame.')]:
        """
        Run object detection on the video. `frame_skip` may be provided here to
        override the constructor default (must be >= 1).
        """
        fs = self.frame_skip if frame_skip is None else frame_skip
        fs = int(fs)

        video_path = video.to_path() if hasattr(video, 'to_path') else str(video)
        cap = cv2.VideoCapture(video_path)
        frame_count = 0
        results_output = []

        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break

            if frame_count % fs == 0:
                # Run object detection on the current frame.
                # The YOLO model returns a list of results for each input image.
                detections = []
                results = self.model(frame)
                for result in results:
                    # Check if result contains boxes.
                    if hasattr(result, 'boxes'):
                        # Prefer per-result class name mapping when available; fallback to model mapping.
                        # Ultralytics may expose this as a dict {id: name} or a list/tuple indexed by class id.
                        names_map = getattr(result, 'names', None) or getattr(self.model, 'names', None) or {}
                        # Get bounding boxes, confidences, and class ids.
                        boxes = result.boxes.xyxy.cpu().numpy() if hasattr(result.boxes.xyxy, 'cpu') else result.boxes.xyxy.numpy()
                        confs = result.boxes.conf.cpu().numpy() if hasattr(result.boxes.conf, 'cpu') else result.boxes.conf.numpy()
                        cls_ids = result.boxes.cls.cpu().numpy() if hasattr(result.boxes.cls, 'cpu') else result.boxes.cls.numpy()
                        for box, conf, cls_id in zip(boxes, confs, cls_ids):
                            cls_int = int(cls_id)
                            if isinstance(names_map, dict):
                                cls_name = names_map.get(cls_int, str(cls_int))
                            elif isinstance(names_map, (list, tuple)):
                                cls_name = names_map[cls_int] if 0 <= cls_int < len(names_map) else str(cls_int)
                            else:
                                cls_name = str(cls_int)
                            detections.append(
                                f"({int(box[0])}, {int(box[1])}, {int(box[2])}, {int(box[3])}) "
                                f"conf: {conf:.2f}, class: {cls_name} ({cls_int})"
                            )

                if detections:
                    results_output.append(f"Frame {frame_count}: " + " | ".join(detections))
                else:
                    results_output.append(f"Frame {frame_count}: (no detections)")
            frame_count += 1

        cap.release()
        return "\n".join(results_output)


class VideoOCR(BaseTool):
    """A tool to recognize optical characters in a video frame by frame.

    Args:
        lang (str | Sequence[str]): The language(s) to recognize.
            Defaults to 'en'.
        frame_skip (int): Process every Nth frame for efficiency.
            Defaults to 1 (process all frames).
        device (str | bool): The device to load the OCR model.
            Defaults to True, meaning automatic selection.
        toolmeta (None | dict | ToolMeta): Additional tool metadata.
            Defaults to None.
    """

    default_desc = 'This tool extracts text from videos using OCR.'

    @require('easyocr')
    def __init__(self,
                 lang: Union[str, Sequence[str]] = 'en',
                 frame_skip: int = 1,
                 device: Union[bool, str] = True,
                 toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        if isinstance(lang, str):
            lang = [lang]
        self.lang = list(lang)
        self.frame_skip = frame_skip
        self.device = device

    def setup(self):
        import easyocr
        self._reader: easyocr.Reader = load_or_build_object(
            easyocr.Reader, self.lang, gpu=self.device)

    def apply(
        self,
        video: str,
        frame_skip: Optional[int] = None,
    ) -> Annotated[str, Info('Extracted text from the video with bounding boxes.')]:
        """
        Extract text from video frames using OCR. `frame_skip` may be provided
        here to override the constructor default (must be >= 1).
        """
        fs = self.frame_skip if frame_skip is None else frame_skip
        fs = int(fs)

        video_path = video.to_path() if hasattr(video, 'to_path') else str(video)
        cap = cv2.VideoCapture(video_path)
        fps = int(cap.get(cv2.CAP_PROP_FPS))
        width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        frame_count = 0

        results = []

        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break

            if frame_count % fs == 0:
                gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
                ocr_results = self._reader.readtext(gray, detail=1)

                frame_text = []
                for bbox, text, _ in ocr_results:
                    bbox_coords = self.extract_bbox(bbox)
                    frame_text.append(f"({bbox_coords[0]}, {bbox_coords[1]}, {bbox_coords[2]}, {bbox_coords[3]}) {text}")
                
                results.append(f"Frame {frame_count}: " + " | ".join(frame_text))

            frame_count += 1

        cap.release()
        return "\n".join(results)

    @staticmethod
    def extract_bbox(bbox) -> Tuple[int, int, int, int]:
        """Convert OCR bbox format to (x1, y1, x2, y2)."""
        xs = [int(point[0]) for point in bbox]
        ys = [int(point[1]) for point in bbox]
        return min(xs), min(ys), max(xs), max(ys)

# class TextToVideoTool(BaseTool):
#     default_desc = (
#         'A reliable video generation tool using Stable Diffusion v1.5. '
#         'Generates videos frame-by-frame with high visual quality and color stability. '
#         'Uses seed-locking technology to ensure temporal consistency.'
#     )

#     def __init__(
#         self,
#         model: str = 'runwayml/stable-diffusion-v1-5',
#         device: str = 'cuda',
#         toolmeta=None,
#     ):
#         super().__init__(toolmeta=toolmeta)
#         self.model_id = model
#         self.device = device
#         self.pipe = None

#     def setup(self):
#         """初始化加载 Stable Diffusion v1.5 模型"""
#         try:
#             # 加载基础模型，使用 float16 以节省显存并加快速度
#             self.pipe = StableDiffusionPipeline.from_pretrained(
#                 self.model_id, 
#                 torch_dtype=torch.float16,
#                 safety_checker=None  # 关闭安全检查以避免不必要的 OOM 或报错
#             )
            
#             # 使用 DPM 调度器，可在 20 步内生成高质量图像
#             self.pipe.scheduler = DPMSolverMultistepScheduler.from_config(
#                 self.pipe.scheduler.config
#             )
            
#             self.pipe.to(self.device)
            
#             # 显存优化技术
#             if self.device == 'cuda':
#                 self.pipe.enable_attention_slicing()
            
#             self.a_prompt = 'best quality, extremely detailed, masterpiece, 4k, cinematic'
#             self.n_prompt = (
#                 'lowres, bad anatomy, bad hands, text, error, missing fingers, '
#                 'extra digit, fewer digits, cropped, worst quality, low quality, '
#                 'jpeg artifacts, signature, watermark, username, blurry'
#             )
#         except Exception as e:
#             print(f"Error loading Stable Diffusion v1.5: {e}")
#             self.pipe = None

#     def apply(
#         self,
#         description: Annotated[str, Info('A text description of the video content.')],
#         output_dir: str = None,
#         num_frames: Annotated[int, Info('Number of frames to generate (e.g. 8-16)')] = 8,
#         resolution: Annotated[str, Info('Resolution as "width,height". Example: "512,512"')] = "512,512",
#         fps: int = 8,
#         filename: str = None,
#     ) -> VideoIO:
        
#         # 1. 解析分辨率
#         try:
#             width, height = [int(p.strip()) for p in resolution.split(',')]
#         except:
#             width, height = 512, 512

#         # 2. 准备输出路径
#         if output_dir is None:
#             output_dir = os.path.join(os.getcwd(), "tools_output")
#         os.makedirs(output_dir, exist_ok=True)

#         if not filename:
#             filename = f"sd_video_{int(time.time())}.mp4"
#         if not filename.endswith('.mp4'):
#             filename += '.mp4'
#         output_path = os.path.join(output_dir, filename)

#         # 3. 如果模型加载失败，生成占位视频
#         if self.pipe is None:
#             print("Warning: Pipe not loaded. Generating fallback video.")
#             self._generate_fallback_video(output_path, width, height, num_frames, fps)
#             return VideoIO(output_path)

#         # 4. 执行逐帧生成逻辑
#         # 核心：锁定随机种子，确保帧与帧之间的内容具有高度相关性
#         generator = torch.Generator(device=self.device).manual_seed(42)
        
#         fourcc = cv2.VideoWriter_fourcc(*'mp4v')
#         video_writer = cv2.VideoWriter(output_path, fourcc, fps, (width, height))

#         print(f"Starting video generation: {num_frames} frames...")
        
        

#         for i in range(num_frames):
#             # 每一帧加入微小的提示词差异，引导模型产生动态感
#             # 例如在描述后加上动态修饰词
#             frame_prompt = f"{description}, {self.a_prompt}, cinematic motion, frame {i+1} of {num_frames}"
            
#             # 直接调用 pipeline，由其内部处理所有精度和 VAE 问题
#             image = self.pipe(
#                 prompt=frame_prompt,
#                 negative_prompt=self.n_prompt,
#                 num_inference_steps=20,
#                 generator=generator,
#                 width=width,
#                 height=height,
#             ).images[0]

#             # PIL (RGB) 转为 OpenCV (BGR)
#             cv_frame = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
#             video_writer.write(cv_frame)
            
#             print(f"Generated frame {i+1}/{num_frames}")

#         video_writer.release()
#         return VideoIO(output_path)

#     def _generate_fallback_video(self, path, w, h, n, fps):
#         """当模型不可用时生成灰色背景视频"""
#         fourcc = cv2.VideoWriter_fourcc(*'mp4v')
#         out = cv2.VideoWriter(path, fourcc, fps, (w, h))
#         for _ in range(n):
#             frame = np.full((h, w, 3), 128, dtype=np.uint8)
#             out.write(frame)
#         out.release()

# class TextToVideoTool(BaseTool):
#     """A tool to generate a video from text description.

#     Args:
#         model (str): The diffusion model to use. Choose from "sd", "sdxl", "sdxl-turbo".
#             Defaults to "sdxl".
#         duration (float): The duration of the generated video in seconds.
#             Defaults to 5.0 seconds.
#         resolution (Tuple[int, int]): The resolution of the generated video (width, height).
#             Defaults to (512, 512).
#         fps (int): Frames per second for the video.
#             Defaults to 24.
#         num_frames (int): Number of frames to generate.
#             Defaults to 24 * 5 (fps * duration).
#         device (str): The device to run the model on. Defaults to "cuda".
#         toolmeta (None | dict | ToolMeta): Additional metadata for the tool.
#     """

#     default_desc = 'This tool generates a video based on a text description using a diffusion model.'

#     @require('diffusers')
#     def __init__(
#         self,
#         model: str = 'sdxl',
#         duration: float = 5.0,
#         resolution: Tuple[int, int] = (512, 512),
#         fps: int = 24,
#         num_frames: int = None,
#         device: str = 'cuda',
#         toolmeta=None,
#     ):
#         super().__init__(toolmeta=toolmeta)
#         assert model in ['sdxl', 'sdxl-turbo']
#         self.model = model
#         self.duration = duration
#         self.resolution = resolution
#         self.fps = fps
#         self.num_frames = num_frames if num_frames else int(fps * duration)
#         self.device = device

#     def setup(self):
#         if self.model == 'sdxl':
#             self.pipe = load_sdxl(device=self.device)
#         elif self.model == 'sdxl-turbo':
#             self.pipe = load_sdxl(model='stabilityai/sdxl-turbo', vae=None, device=self.device)
#         self.a_prompt = 'best quality, extremely detailed'
#         self.n_prompt = 'longbody, lowres, bad anatomy, bad hands, '\
#                         'missing fingers, extra digit, fewer digits, '\
#                         'cropped, worst quality, low quality'

#     def apply(
#         self,
#         description: Annotated[str, Info('A text description of the video content.')],
#         output_dir: str = None,
#         duration: Optional[float] = None,
#         resolution: Annotated[Optional[str], Info('Resolution as "width,height". Example: "512,512"')] = None,
#         fps: Optional[int] = None,
#         num_frames: Optional[int] = None,
#         filename: Optional[str] = None,
#     ) -> VideoIO:
#         """
#         Generate a video from `description`.

#         Call-time overrides are supported for: duration, resolution, fps, num_frames and filename.
#         If an override is omitted, the instance defaults from the constructor are used.
#         """
#         # Parse resolution string to tuple if provided
#         res_tuple = None
#         if resolution is not None:
#             try:
#                 parts = [int(float(p.strip())) for p in resolution.split(',')]
#                 if len(parts) != 2:
#                     raise ValueError
#                 res_tuple = (parts[0], parts[1])
#             except Exception:
#                 raise ValueError('resolution must be a string in the format "width,height", e.g., "512,512"')

#         output_video = self.generate_video_from_text(
#             description,
#             output_dir=output_dir,
#             duration=duration,
#             resolution=res_tuple,
#             fps=fps,
#             num_frames=num_frames,
#             filename=filename,
#         )
#         return VideoIO(output_video)

#     def generate_video_from_text(
#         self,
#         description: str,
#         output_dir: str = None,
#         duration: Optional[float] = None,
#         resolution: Optional[Tuple[int, int]] = None,
#         fps: Optional[int] = None,
#         num_frames: Optional[int] = None,
#         filename: Optional[str] = None,
#     ) -> str:
#         """
#         Internal helper that generates a video from text. All parameters may be
#         overridden at call time; any omitted parameter falls back to the
#         instance defaults configured in the constructor.
#         """
#         # Resolve runtime overrides -> final values
#         width, height = resolution if resolution is not None else self.resolution
#         fps = int(fps) if fps is not None else int(self.fps)

#         if num_frames is not None:
#             n_frames = int(num_frames)
#         else:
#             # If duration provided, compute frames from duration; otherwise use instance num_frames
#             if duration is not None:
#                 n_frames = int(fps * duration)
#             else:
#                 n_frames = int(self.num_frames)

#         import os, time
#         if output_dir is None:
#             output_dir = os.path.join(os.getcwd(), "tools_output")
#         os.makedirs(output_dir, exist_ok=True)

#         # Build a safe filename
#         if filename:
#             safe_name = os.path.basename(filename)
#             if not safe_name.lower().endswith('.mp4'):
#                 safe_name = f"{safe_name}.mp4"
#         else:
#             # include a short sanitized description and a timestamp to avoid collisions
#             desc_snip = ''.join(c for c in description[:20] if c.isalnum() or c in (' ', '_')).replace(' ', '_')
#             safe_name = f"text_to_video_{desc_snip}_{int(time.time())}.mp4"

#         output_path = os.path.join(output_dir, safe_name)
#         fourcc = cv2.VideoWriter_fourcc(*'mp4v')
#         out = cv2.VideoWriter(output_path, fourcc, fps, (width, height))

#         # Generate multiple frames with slight variations
#         for i in range(n_frames):
#             prompt = f'{description}, frame {i+1}, {self.a_prompt}'
#             # If self.pipe is not available (lite mode), skip actual generation to avoid crashing
#             if self.pipe is None:
#                 # create a blank frame (gray) as a fallback
#                 frame = np.full((height, width, 3), 128, dtype=np.uint8)
#             else:
#                 image = self.pipe(prompt, num_inference_steps=30, negative_prompt=self.n_prompt).images[0]
#                 frame = np.array(image)[:, :, ::-1]  # Convert RGB to BGR for OpenCV
#             out.write(frame)

#         out.release()
#         return output_path

class TextToVideoTool(BaseTool):
    default_desc = (
        'Generate a simple MP4 video without heavy dependencies. '
        'Uses PIL to synthesize frames and imageio/cv2 to encode.'
    )

    def __init__(
        self,
        model: str = 'lite-synthetic-generator',
        device: str = 'cpu',
        toolmeta=None,
    ):
        super().__init__(toolmeta=toolmeta)
        self.model_id = model
        self.device = device
        self.pipe = None
        # Prompts kept for compatibility with prior tests
        self.a_prompt = 'best quality, extremely detailed, masterpiece, cinematic'
        self.n_prompt = 'lowres, artifacts, blurry'

    def setup(self):
        # Lite setup: no heavy model; keep attributes initialized
        self.pipe = None

    def _parse_resolution(self, resolution: str) -> Tuple[int, int]:
        try:
            parts = resolution.split(',')
            w = int(float(parts[0].strip()))
            h = int(float(parts[1].strip()))
            return max(16, w), max(16, h)
        except Exception:
            return 128, 128

    def _generate_frame(self, w: int, h: int, t: int, T: int, text: str):
        from PIL import Image, ImageDraw, ImageFont
        img = Image.new('RGB', (w, h), (240, 240, 240))
        draw = ImageDraw.Draw(img)

        # Moving square
        margin = max(4, min(w, h) // 16)
        sq = max(8, min(w, h) // 6)
        x = margin + int((w - 2 * margin - sq) * (t / max(1, T - 1)))
        y = margin + int((h - 2 * margin - sq) * (0.5 - 0.5 * (t / max(1, T - 1))))
        color = (
            50 + (t * 30) % 205,
            80 + (t * 50) % 175,
            110 + (t * 70) % 145,
        )
        draw.rectangle([x, y, x + sq, y + sq], fill=color)

        # Text overlay
        overlay = f"{text[:32]}"
        try:
            font = ImageFont.load_default()
        except Exception:
            font = None
        draw.text((margin, h - margin - 12), overlay, fill=(30, 30, 30), font=font)
        return img

    def _write_video_imageio(self, frames, path: str, fps: int) -> bool:
        try:
            import imageio
            # Prefer libx264 if available; fallback handled by imageio
            writer = imageio.get_writer(path, fps=fps)
            for fr in frames:
                writer.append_data(fr)
            writer.close()
            return True
        except Exception as e:
            print(f"imageio writer failed: {e}")
            return False

    def _write_video_cv2(self, frames, path: str, fps: int) -> bool:
        try:
            import cv2
            h, w, _ = frames[0].shape
            fourcc = cv2.VideoWriter_fourcc(*'mp4v')
            out = cv2.VideoWriter(path, fourcc, fps, (w, h))
            for fr in frames:
                out.write(cv2.cvtColor(fr, cv2.COLOR_RGB2BGR))
            out.release()
            return True
        except Exception as e:
            print(f"cv2 writer failed: {e}")
            return False

    def apply(
        self,
        description: Annotated[str, Info('A text description of the video content.')],
        output_dir: str = None,
        num_frames: Annotated[int, Info('Number of frames to generate (e.g. 8-16)')] = 8,
        resolution: Annotated[str, Info('Resolution as "width,height". Example: "512,512"')] = "128,128",
        fps: int = 8,
        filename: str = None,
    ) -> VideoIO:
        import os
        import numpy as np
        from datetime import datetime

        # Resolve output directory
        if output_dir is None:
            output_dir = DEFAULT_TOOLS_OUTPUT_DIR
        os.makedirs(output_dir, exist_ok=True)

        # Resolve filename
        if not filename:
            ts = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"text_to_video_{ts}"

        # Parse resolution
        w, h = self._parse_resolution(resolution)

        # Generate frames
        frames = []
        for t in range(max(1, num_frames)):
            img = self._generate_frame(w, h, t, max(1, num_frames), description)
            frames.append(np.array(img))

        # Write MP4 via imageio first, then fallback to cv2
        mp4_path = os.path.join(output_dir, f"{filename}.mp4")
        ok = self._write_video_imageio(frames, mp4_path, max(1, fps))
        if not ok:
            ok = self._write_video_cv2(frames, mp4_path, max(1, fps))

        if not ok:
            raise RuntimeError(
                "Failed to encode MP4. Please install 'imageio[ffmpeg]' or 'opencv-python'."
            )

        return VideoIO(mp4_path)

class AddTextVideoTool(BaseTool):
    """一个用于向视频中添加文字的工具。

    参数:
        text (str): 要添加到视频中的文字内容。
        position (Tuple[int, int]): 文字在视频帧中的位置 (x, y)，单位为像素，默认为 (50, 50)。
        font (int): OpenCV 中使用的字体类型，默认为 cv2.FONT_HERSHEY_SIMPLEX。
        font_scale (float): 文字的缩放因子，默认为 1.0。
        color (Tuple[int, int, int]): 文字颜色（B, G, R），默认为白色 (255, 255, 255)。
        thickness (int): 文字的粗细，默认为 2。
        toolmeta: 其他工具元数据，默认为 None。
    """

    default_desc = "此工具用于在视频中添加文字水印或说明。"

    def __init__(self,
                 text: str = "Hello, World!",
                 position: Tuple[int, int] = (50, 50),
                 font: int = cv2.FONT_HERSHEY_SIMPLEX,
                 font_scale: float = 1.0,
                 color: Tuple[int, int, int] = (255, 255, 255),
                 thickness: int = 2,
                 toolmeta=None):
        super().__init__(toolmeta=toolmeta)
        self.text = text
        self.position = position
        self.font = font
        self.font_scale = font_scale
        self.color = color
        self.thickness = thickness

    def setup(self):
        # 不需要额外的初始化设置
        pass

    def apply(
        self,
        video: str,
        output_dir: str = None,
        text: Optional[str] = None,
        position: Annotated[Optional[str], Info('Text position as "x,y" in pixels, e.g., "50,50"')] = None,
        font: Optional[int] = None,
        font_scale: Optional[float] = None,
        color: Annotated[Optional[str], Info('Text color as "B,G,R" (0-255), e.g., "255,255,255"')] = None,
        thickness: Optional[int] = None,
        filename: Optional[str] = None,
    ) -> Annotated[VideoIO, Info("Video with text overlay")]:
        """
        Overlay text on the input `video` and write a new video file.

        Call-time overrides (text, position, font, font_scale, color, thickness)
        are accepted; if omitted, the instance defaults from __init__ are used.
        `filename` may be provided to control output basename; otherwise a
        timestamped name is created under `tools_output`.
        """
        # Resolve overrides
        txt = self.text if text is None else text
        pos = self.position if position is None else position
        fnt = self.font if font is None else font
        fscale = self.font_scale if font_scale is None else font_scale
        col = self.color if color is None else color
        thick = self.thickness if thickness is None else thickness

        # Basic validation and parsing
        # position: allow string "x,y" or a tuple/list
        if isinstance(pos, str):
            try:
                p = [int(float(v.strip())) for v in pos.split(',')]
                if len(p) != 2:
                    raise ValueError
                pos = (p[0], p[1])
            except Exception:
                raise ValueError('position must be a string in the format "x,y", e.g., "50,50"')
        elif isinstance(pos, (tuple, list)) and len(pos) == 2:
            try:
                pos = (int(pos[0]), int(pos[1]))
            except Exception:
                raise ValueError("position values must be convertible to int")
        else:
            raise ValueError("position must be provided as a string 'x,y' or a tuple/list (x, y)")

        # color: allow string "B,G,R" or a tuple/list of length 3
        if isinstance(col, str):
            try:
                c = [int(float(v.strip())) for v in col.split(',')]
                if len(c) != 3:
                    raise ValueError
                col = (c[0], c[1], c[2])
            except Exception:
                raise ValueError('color must be a string in the format "B,G,R", e.g., "255,255,255"')
        elif isinstance(col, (tuple, list)) and len(col) == 3:
            try:
                col = (int(col[0]), int(col[1]), int(col[2]))
            except Exception:
                raise ValueError("color values must be convertible to int")
        else:
            raise ValueError("color must be provided as a string 'B,G,R' or a tuple/list (B, G, R)")
        for v in col:
            if v < 0 or v > 255:
                raise ValueError("color values must be in range [0, 255]")

        if fscale is None:
            fscale = 1.0
        else:
            fscale = float(fscale)
            if fscale <= 0:
                raise ValueError("font_scale must be > 0")

        if thick is None:
            thick = 1
        else:
            thick = int(thick)
            if thick < 0:
                raise ValueError("thickness must be >= 0")

        # Input is expected to be a file path string.
        video_path = video.to_path() if hasattr(video, 'to_path') else str(video)
        cap = cv2.VideoCapture(video_path)
        if not cap.isOpened():
            raise ValueError("无法打开视频文件: " + video_path)

        # Get video properties
        fps = cap.get(cv2.CAP_PROP_FPS)
        width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))

        # Prepare output path
        import os, time
        if output_dir is None:
            output_dir = DEFAULT_TOOLS_OUTPUT_DIR
        os.makedirs(output_dir, exist_ok=True)

        if filename:
            out_name = os.path.basename(filename)
            if not out_name.lower().endswith('.mp4'):
                out_name = f"{out_name}.mp4"
        else:
            desc_snip = ''.join(c for c in (txt or "")[:20] if c.isalnum() or c in (' ', '_')).replace(' ', '_')
            out_name = f"video_with_text_{desc_snip}_{int(time.time())}.mp4"

        output_path = os.path.join(output_dir, out_name)

        fourcc = cv2.VideoWriter_fourcc(*'mp4v')
        out = cv2.VideoWriter(output_path, fourcc, fps, (width, height))

        # Process frames and overlay text
        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break

            cv2.putText(frame,
                        str(txt),
                        pos,
                        fnt,
                        fscale,
                        col,
                        thick,
                        cv2.LINE_AA)
            out.write(frame)

        cap.release()
        out.release()
        return VideoIO(output_path)

